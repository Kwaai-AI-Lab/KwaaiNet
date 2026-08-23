"""Apply Kwaai branding to a pandoc-generated pptx.

Colours sampled from the logo in the existing intern deck:
  blue  #0A63DC (cloud + wordmark)   deep #0A3FB4   brown #4A3520 (tree)
"""
import sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor

BLUE  = RGBColor(0x0A, 0x63, 0xDC)
DEEP  = RGBColor(0x0A, 0x3F, 0xB4)
BROWN = RGBColor(0x4A, 0x35, 0x20)
INK   = RGBColor(0x1A, 0x1A, 0x1A)

src, logo, dst = sys.argv[1], sys.argv[2], sys.argv[3]
p = Presentation(src)
SW, SH = p.slide_width, p.slide_height

def title_of(slide):
    for sh in slide.shapes:
        if sh.is_placeholder and sh.placeholder_format.idx == 0 and sh.has_text_frame:
            return sh
    for sh in slide.shapes:                      # fall back to first text shape
        if sh.has_text_frame and sh.text.strip():
            return sh
    return None

def colour(shape, rgb, bold=True, size=None):
    if shape is None: return
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = rgb
            run.font.bold = bold
            if size: run.font.size = Pt(size)

n_logo = n_rule = 0
for i, slide in enumerate(p.slides):
    layout = slide.slide_layout.name
    t = title_of(slide)

    if i == 0:
        # Title slide: large logo, top-centre, title pushed clear of it.
        h = Inches(1.35)
        w = int(h * 777 / 800)
        slide.shapes.add_picture(logo, int((SW - w) / 2), Inches(0.22), width=w, height=h)
        if t is not None:
            t.top = Inches(1.75)
            colour(t, DEEP, bold=True)
        for sh in slide.shapes:
            if sh.is_placeholder and sh.placeholder_format.idx != 0 and sh.has_text_frame:
                # only nudge down where it still fits; the date placeholder sits
                # near the foot already and would otherwise land off-canvas
                if int(sh.top) + int(sh.height) + Inches(0.55) < SH:
                    sh.top = Emu(int(sh.top) + Inches(0.55))
                colour(sh, BROWN, bold=False)
        n_logo += 1
        continue

    # Every other slide: small logo, bottom-right. Top-right collides with
    # pandoc's content placeholders and tables, which start at y≈0.22in.
    h = Inches(0.40)
    w = int(h * 777 / 800)
    slide.shapes.add_picture(logo, SW - w - Inches(0.20), SH - h - Inches(0.12),
                             width=w, height=h)
    n_logo += 1

    if layout == "Section Header":
        colour(t, DEEP, bold=True)
        # brown rule under the week heading, positioned from the title itself
        ry = (int(t.top) + int(t.height) + Inches(0.06)) if t is not None else Inches(4.7)
        ry = min(ry, SH - Inches(0.30))
        rule = slide.shapes.add_shape(1, Inches(0.55), ry,
                                      SW - Inches(1.10), Emu(28575))  # ~0.03"
        rule.fill.solid(); rule.fill.fore_color.rgb = BROWN
        rule.line.fill.background(); rule.shadow.inherit = False
        n_rule += 1
    else:
        colour(t, BLUE, bold=True)

p.save(dst)
print(f"  branded: {len(p.slides)} slides, {n_logo} logos, {n_rule} section rules")
